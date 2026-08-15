from flask import Flask, jsonify, request, send_from_directory
from flask_cors import CORS
from datetime import datetime, timezone
import os
import sqlite3
import json
import re
import uuid
from dotenv import load_dotenv
import PyPDF2
from pptx import Presentation
from werkzeug.utils import secure_filename

load_dotenv()

app = Flask(__name__)

# Configuration
app.config['SECRET_KEY'] = os.getenv('SECRET_KEY', 'dev-secret-key')
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024
app.config['UPLOAD_FOLDER'] = os.path.join(os.path.dirname(os.path.abspath(__file__)), '..', 'uploads')
app.config['ALLOWED_EXTENSIONS'] = {'pdf', 'pptx', 'ppt', 'docx', 'txt'}

os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

CORS(app)

# ==================== DATABASE SETUP ====================

def get_db():
    db_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), '..', 'data', 'notes.db')
    os.makedirs(os.path.dirname(db_path), exist_ok=True)
    conn = sqlite3.connect(db_path)
    conn.row_factory = sqlite3.Row
    return conn

def init_db():
    conn = get_db()
    cursor = conn.cursor()
    
    # Notes table
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS notes (
            note_id TEXT PRIMARY KEY,
            title TEXT NOT NULL,
            course_code TEXT,
            description TEXT,
            summary TEXT,
            key_points TEXT,
            word_count INTEGER,
            file_name TEXT,
            file_url TEXT,
            tags TEXT,
            upload_date TEXT,
            views INTEGER DEFAULT 0,
            downloads INTEGER DEFAULT 0,
            rating REAL DEFAULT 0,
            rating_count INTEGER DEFAULT 0,
            collection TEXT
        )
    ''')
    
    # Quizzes table
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS quizzes (
            quiz_id TEXT PRIMARY KEY,
            title TEXT NOT NULL,
            course_code TEXT,
            description TEXT,
            difficulty TEXT DEFAULT 'medium',
            questions TEXT,
            question_count INTEGER DEFAULT 0,
            file_name TEXT,
            file_url TEXT,
            tags TEXT,
            upload_date TEXT,
            views INTEGER DEFAULT 0,
            downloads INTEGER DEFAULT 0,
            rating REAL DEFAULT 0,
            rating_count INTEGER DEFAULT 0,
            collection TEXT
        )
    ''')
    
    # Quiz Attempts table
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS quiz_attempts (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            quiz_id TEXT,
            student_name TEXT,
            score INTEGER,
            total_questions INTEGER,
            percentage REAL,
            passed BOOLEAN,
            answers TEXT,
            completed_at TEXT,
            FOREIGN KEY (quiz_id) REFERENCES quizzes (quiz_id)
        )
    ''')
    
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS ratings (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            note_id TEXT,
            rating INTEGER,
            user_id TEXT,
            created_at TEXT,
            FOREIGN KEY (note_id) REFERENCES notes (note_id)
        )
    ''')
    
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS collections (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            name TEXT UNIQUE,
            created_at TEXT,
            description TEXT
        )
    ''')
    
    conn.commit()
    conn.close()
    print("✅ SQLite database initialized!")

init_db()

# ==================== TEXT EXTRACTION ====================

def extract_text_from_pdf(file_path):
    content = ""
    try:
        with open(file_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    content += text + "\n"
        return content
    except Exception as e:
        print(f"PDF error: {e}")
        return ""

def extract_text_from_pptx(file_path):
    content = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    content += shape.text + "\n"
        return content
    except Exception as e:
        print(f"PPTX error: {e}")
        return ""

# ==================== SMART AI PROCESSOR ====================

class SmartAIProcessor:
    def __init__(self):
        self.provider = None
        self.client = None
        
        groq_key = os.getenv('GROQ_API_KEY')
        if groq_key:
            try:
                from groq import Groq
                self.client = Groq(api_key=groq_key)
                self.provider = 'groq'
                self.model = 'llama-3.3-70b-versatile'
                print("✅ Groq AI connected!")
                return
            except Exception as e:
                print(f"⚠️ Groq error: {e}")
        
        gemini_key = os.getenv('GEMINI_API_KEY')
        if gemini_key:
            try:
                import google.generativeai as genai
                genai.configure(api_key=gemini_key)
                self.client = genai
                self.provider = 'gemini'
                self.model = 'gemini-1.5-flash'
                print("✅ Gemini AI connected!")
                return
            except Exception as e:
                print(f"⚠️ Gemini error: {e}")
        
        if not self.provider:
            print("⚠️ No AI provider found. Using fallback text processing.")
    
    def generate_notes(self, text, title=""):
        if not text or len(text) < 100:
            return {
                'summary': 'Insufficient text to generate notes.',
                'key_points': ['No key points could be extracted.'],
                'title': title or 'Untitled Note',
                'word_count': len(text.split()) if text else 0
            }
        
        text = re.sub(r'\s+', ' ', text)
        text = text[:6000]
        
        if not self.provider:
            return self._fallback_notes(text, title)
        
        prompt = f"""You are an expert study note creator. Create comprehensive, detailed study notes from the following content.

CONTENT:
{text}

Create a JSON response with:

1. "summary": A DETAILED paragraph summary (100-150 words)
2. "key_points": A list of 6-10 IMPORTANT key points as complete sentences

Return ONLY valid JSON with this structure:
{{"summary": "Your detailed summary here.", "key_points": ["Point one.", "Point two.", "Point three."]}}"""

        try:
            if self.provider == 'groq':
                response = self.client.chat.completions.create(
                    model=self.model,
                    messages=[
                        {"role": "system", "content": "You are an expert study note creator. Return ONLY valid JSON."},
                        {"role": "user", "content": prompt}
                    ],
                    temperature=0.5,
                    max_tokens=1200
                )
                result_text = response.choices[0].message.content
            elif self.provider == 'gemini':
                model = self.client.GenerativeModel(self.model)
                response = model.generate_content(prompt)
                result_text = response.text
            
            result_text = result_text.strip()
            if '```json' in result_text:
                result_text = result_text.split('```json')[1].split('```')[0]
            elif '```' in result_text:
                result_text = result_text.split('```')[1].split('```')[0]
            
            json_match = re.search(r'\{[\s\S]*\}', result_text)
            if json_match:
                result_text = json_match.group()
            
            data = json.loads(result_text)
            summary = data.get('summary', 'No summary generated.')
            key_points = data.get('key_points', ['No key points extracted.'])
            
            if not isinstance(key_points, list):
                key_points = [str(key_points)]
            key_points = [str(k).strip() for k in key_points if k and str(k).strip()]
            if not key_points:
                key_points = ['No key points extracted.']
            
            return {
                'summary': summary,
                'key_points': key_points[:10],
                'title': title or 'Generated Notes',
                'word_count': len(text.split())
            }
        except Exception as e:
            print(f"❌ AI generation error: {e}")
            return self._fallback_notes(text, title)
    
    def generate_quiz_questions(self, text, num_questions=10, difficulty="medium"):
        """Generate quiz questions from content"""
        if not text or len(text) < 100:
            return self._fallback_quiz_questions(num_questions)
        
        text = re.sub(r'\s+', ' ', text)
        text = text[:6000]
        
        if not self.provider:
            return self._fallback_quiz_questions(num_questions)
        
        diff_instruction = {
            'easy': 'Make questions straightforward and focus on basic concepts.',
            'medium': 'Make questions moderate difficulty testing understanding.',
            'hard': 'Make challenging questions that require deep understanding and analysis.'
        }.get(difficulty, 'Make questions moderate difficulty.')
        
        prompt = f"""You are an expert quiz creator. Generate {num_questions} multiple-choice questions from the content below.

CONTENT:
{text}

DIFFICULTY LEVEL: {difficulty}
{diff_instruction}

Each question must have exactly 4 options with ONE correct answer.

Return ONLY valid JSON with this structure:
{{
    "questions": [
        {{
            "question": "Question text?",
            "options": ["Option A", "Option B", "Option C", "Option D"],
            "correct_answer": 0,
            "explanation": "Brief explanation of why this is correct."
        }}
    ]
}}"""

        try:
            if self.provider == 'groq':
                response = self.client.chat.completions.create(
                    model=self.model,
                    messages=[
                        {"role": "system", "content": "You are an expert quiz creator. Return ONLY valid JSON."},
                        {"role": "user", "content": prompt}
                    ],
                    temperature=0.4,
                    max_tokens=1500
                )
                result_text = response.choices[0].message.content
            elif self.provider == 'gemini':
                model = self.client.GenerativeModel(self.model)
                response = model.generate_content(prompt)
                result_text = response.text
            
            result_text = result_text.strip()
            if '```json' in result_text:
                result_text = result_text.split('```json')[1].split('```')[0]
            elif '```' in result_text:
                result_text = result_text.split('```')[1].split('```')[0]
            
            json_match = re.search(r'\{[\s\S]*\}', result_text)
            if json_match:
                result_text = json_match.group()
            
            data = json.loads(result_text)
            questions = data.get('questions', [])
            
            valid_questions = []
            for q in questions:
                if 'question' in q and 'options' in q and len(q['options']) >= 4 and 'correct_answer' in q:
                    valid_questions.append(q)
            
            if valid_questions:
                return valid_questions
            return self._fallback_quiz_questions(num_questions)
        except Exception as e:
            print(f"❌ Quiz generation error: {e}")
            return self._fallback_quiz_questions(num_questions)
    
    def _fallback_quiz_questions(self, num_questions):
        questions = []
        for i in range(min(num_questions, 10)):
            questions.append({
                'question': f"Sample Question {i+1}: What is a key concept from this content?",
                'options': [
                    "Understanding core principles is essential.",
                    "Memorizing facts is the main goal.",
                    "Practical examples are the focus.",
                    "Historical context is most important."
                ],
                'correct_answer': 0,
                'explanation': "Understanding core principles is fundamental to mastering any subject."
            })
        return questions
    
    def _fallback_notes(self, text, title=""):
        text = re.sub(r'\s+', ' ', text)
        sentences = re.split(r'[.!?]+\s+', text)
        sentences = [s.strip() for s in sentences if len(s.strip()) > 20]
        
        if sentences:
            summary = '. '.join(sentences[:2]) + '.'
            if len(summary) > 300:
                summary = summary[:300] + '...'
            key_points = sentences[2:7] if len(sentences) >= 7 else sentences[2:]
            if not key_points:
                key_points = sentences[:5]
        else:
            summary = "No summary could be generated."
            key_points = ["No key points could be extracted."]
        
        return {
            'summary': summary,
            'key_points': key_points,
            'title': title or 'Untitled Note',
            'word_count': len(text.split())
        }

ai_processor = SmartAIProcessor()

# ==================== ROUTES ====================

@app.route('/')
def home():
    return jsonify({
        'message': 'USP Notes Marketplace API is running!',
        'status': 'active',
        'database': 'SQLite',
        'ai_model': ai_processor.provider or 'Fallback'
    })

# ===== NOTES ROUTES =====

@app.route('/api/notes', methods=['GET'])
def get_all_notes():
    try:
        conn = get_db()
        cursor = conn.cursor()
        cursor.execute('SELECT * FROM notes ORDER BY upload_date DESC')
        rows = cursor.fetchall()
        conn.close()
        
        notes = []
        for row in rows:
            note = dict(row)
            note['key_points'] = json.loads(note['key_points']) if note['key_points'] else []
            note['tags'] = json.loads(note['tags']) if note['tags'] else []
            notes.append(note)
        
        return jsonify(notes), 200
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/notes/<note_id>', methods=['GET'])
def get_note(note_id):
    try:
        conn = get_db()
        cursor = conn.cursor()
        cursor.execute('SELECT * FROM notes WHERE note_id = ?', (note_id,))
        row = cursor.fetchone()
        conn.close()
        
        if not row:
            return jsonify({'error': 'Note not found'}), 404
        
        note = dict(row)
        note['key_points'] = json.loads(note['key_points']) if note['key_points'] else []
        note['tags'] = json.loads(note['tags']) if note['tags'] else []
        return jsonify(note), 200
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/notes/search', methods=['GET'])
def search_notes():
    try:
        query = request.args.get('q', '')
        if not query:
            return jsonify([]), 200
        
        conn = get_db()
        cursor = conn.cursor()
        cursor.execute('''
            SELECT * FROM notes 
            WHERE title LIKE ? OR course_code LIKE ? OR tags LIKE ? OR summary LIKE ?
            ORDER BY upload_date DESC
        ''', (f'%{query}%', f'%{query}%', f'%{query}%', f'%{query}%'))
        
        rows = cursor.fetchall()
        conn.close()
        
        notes = []
        for row in rows:
            note = dict(row)
            note['key_points'] = json.loads(note['key_points']) if note['key_points'] else []
            note['tags'] = json.loads(note['tags']) if note['tags'] else []
            notes.append(note)
        
        return jsonify(notes), 200
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/collections', methods=['GET'])
def get_collections():
    try:
        conn = get_db()
        cursor = conn.cursor()
        cursor.execute('''
            SELECT collection, COUNT(*) as count 
            FROM notes 
            WHERE collection IS NOT NULL AND collection != ''
            GROUP BY collection
            ORDER BY collection
        ''')
        rows = cursor.fetchall()
        conn.close()
        
        collections = []
        for row in rows:
            collections.append({
                'name': row['collection'],
                'count': row['count']
            })
        
        return jsonify(collections), 200
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/collections/<collection_name>', methods=['GET'])
def get_collection_notes(collection_name):
    try:
        conn = get_db()
        cursor = conn.cursor()
        cursor.execute('''
            SELECT * FROM notes WHERE collection = ?
            ORDER BY upload_date DESC
        ''', (collection_name,))
        rows = cursor.fetchall()
        conn.close()
        
        notes = []
        for row in rows:
            note = dict(row)
            note['key_points'] = json.loads(note['key_points']) if note['key_points'] else []
            note['tags'] = json.loads(note['tags']) if note['tags'] else []
            notes.append(note)
        
        return jsonify(notes), 200
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/upload', methods=['POST'])
def upload_note():
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file uploaded'}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400
        
        if not allowed_file(file.filename):
            return jsonify({'error': 'File type not allowed'}), 400
        
        title = request.form.get('title', file.filename.rsplit('.', 1)[0])
        course_code = request.form.get('course_code', '').upper()
        description = request.form.get('description', '')
        tags_input = request.form.get('tags', '')
        tags = [t.strip() for t in tags_input.split(',') if t.strip()]
        collection = request.form.get('collection', '').strip()
        
        filename = secure_filename(file.filename)
        os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(file_path)
        
        content = ""
        if filename.endswith('.pdf'):
            content = extract_text_from_pdf(file_path)
        elif filename.endswith(('.pptx', '.ppt')):
            content = extract_text_from_pptx(file_path)
        elif filename.endswith('.txt'):
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
        else:
            return jsonify({'error': 'Could not extract text from this file type'}), 400
        
        print(f"📄 Extracted {len(content)} characters")
        
        note_data = ai_processor.generate_notes(content, title)
        
        note_id = str(uuid.uuid4())[:8]
        
        conn = get_db()
        cursor = conn.cursor()
        
        if collection:
            cursor.execute('INSERT OR IGNORE INTO collections (name, created_at) VALUES (?, ?)',
                          (collection, datetime.now(timezone.utc).isoformat()))
        
        cursor.execute('''
            INSERT INTO notes (
                note_id, title, course_code, description, summary, key_points,
                word_count, file_name, file_url, tags, upload_date, 
                views, downloads, rating, rating_count, collection
            ) VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
        ''', (
            note_id,
            note_data['title'],
            course_code,
            description,
            note_data['summary'],
            json.dumps(note_data['key_points']),
            note_data['word_count'],
            filename,
            f"/uploads/{filename}",
            json.dumps(tags),
            datetime.now(timezone.utc).isoformat(),
            0, 0, 0.0, 0,
            collection
        ))
        conn.commit()
        conn.close()
        
        return jsonify({
            'success': True,
            'message': 'Note uploaded and generated successfully!',
            'note': {
                'note_id': note_id,
                'title': note_data['title'],
                'summary': note_data['summary'],
                'key_points': note_data['key_points'],
                'word_count': note_data['word_count'],
                'course_code': course_code,
                'description': description,
                'tags': tags,
                'collection': collection
            }
        }), 201
        
    except Exception as e:
        print(f"❌ Upload error: {e}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/notes/<note_id>/rate', methods=['POST'])
def rate_note(note_id):
    try:
        data = request.get_json()
        rating = data.get('rating', 0)
        
        if not 1 <= rating <= 5:
            return jsonify({'error': 'Rating must be between 1 and 5'}), 400
        
        conn = get_db()
        cursor = conn.cursor()
        cursor.execute('SELECT rating, rating_count FROM notes WHERE note_id = ?', (note_id,))
        row = cursor.fetchone()
        
        if not row:
            conn.close()
            return jsonify({'error': 'Note not found'}), 404
        
        current_rating = row['rating'] or 0
        current_count = row['rating_count'] or 0
        
        new_count = current_count + 1
        new_rating = round(((current_rating * current_count) + rating) / new_count, 1)
        
        cursor.execute('''
            UPDATE notes SET rating = ?, rating_count = ? WHERE note_id = ?
        ''', (new_rating, new_count, note_id))
        
        conn.commit()
        conn.close()
        
        return jsonify({
            'success': True,
            'rating': new_rating,
            'rating_count': new_count
        }), 200
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/notes/<note_id>/view', methods=['POST'])
def increment_views(note_id):
    try:
        conn = get_db()
        cursor = conn.cursor()
        cursor.execute('UPDATE notes SET views = views + 1 WHERE note_id = ?', (note_id,))
        conn.commit()
        conn.close()
        return jsonify({'success': True}), 200
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/notes/<note_id>/download', methods=['POST'])
def increment_downloads(note_id):
    try:
        conn = get_db()
        cursor = conn.cursor()
        cursor.execute('UPDATE notes SET downloads = downloads + 1 WHERE note_id = ?', (note_id,))
        conn.commit()
        conn.close()
        return jsonify({'success': True}), 200
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# ===== QUIZ ROUTES =====

@app.route('/api/quizzes', methods=['GET'])
def get_all_quizzes():
    try:
        conn = get_db()
        cursor = conn.cursor()
        cursor.execute('SELECT * FROM quizzes ORDER BY upload_date DESC')
        rows = cursor.fetchall()
        conn.close()
        
        quizzes = []
        for row in rows:
            quiz = dict(row)
            quiz['questions'] = json.loads(quiz['questions']) if quiz['questions'] else []
            quiz['tags'] = json.loads(quiz['tags']) if quiz['tags'] else []
            quiz['question_count'] = len(quiz['questions'])
            quizzes.append(quiz)
        
        return jsonify(quizzes), 200
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/quizzes/<quiz_id>', methods=['GET'])
def get_quiz(quiz_id):
    try:
        conn = get_db()
        cursor = conn.cursor()
        cursor.execute('SELECT * FROM quizzes WHERE quiz_id = ?', (quiz_id,))
        row = cursor.fetchone()
        conn.close()
        
        if not row:
            return jsonify({'error': 'Quiz not found'}), 404
        
        quiz = dict(row)
        quiz['questions'] = json.loads(quiz['questions']) if quiz['questions'] else []
        quiz['tags'] = json.loads(quiz['tags']) if quiz['tags'] else []
        
        # Remove correct answers for students
        for q in quiz['questions']:
            q.pop('correct_answer', None)
        
        return jsonify(quiz), 200
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/quizzes/<quiz_id>/submit', methods=['POST'])
def submit_quiz(quiz_id):
    try:
        data = request.get_json()
        student_name = data.get('student_name', 'Anonymous')
        user_answers = data.get('answers', {})
        
        conn = get_db()
        cursor = conn.cursor()
        cursor.execute('SELECT * FROM quizzes WHERE quiz_id = ?', (quiz_id,))
        row = cursor.fetchone()
        conn.close()
        
        if not row:
            return jsonify({'error': 'Quiz not found'}), 404
        
        quiz = dict(row)
        questions = json.loads(quiz['questions']) if quiz['questions'] else []
        
        score = 0
        results = []
        
        for i, question in enumerate(questions):
            user_ans = user_answers.get(str(i))
            correct_idx = question.get('correct_answer', 0)
            
            is_correct = False
            if user_ans is not None:
                try:
                    user_ans_int = int(user_ans)
                    is_correct = (user_ans_int == correct_idx)
                    if is_correct:
                        score += 1
                except:
                    pass
            
            options = question.get('options', ['A', 'B', 'C', 'D'])
            user_text = options[user_ans] if user_ans is not None and isinstance(user_ans, int) and user_ans < len(options) else 'Not answered'
            correct_text = options[correct_idx] if isinstance(correct_idx, int) and correct_idx < len(options) else 'Unknown'
            
            results.append({
                'question': question.get('question', ''),
                'your_answer': user_text,
                'correct_answer': correct_text,
                'is_correct': is_correct,
                'explanation': question.get('explanation', 'No explanation provided.')
            })
        
        total = len(questions)
        percentage = round((score / total) * 100, 2) if total > 0 else 0
        passed = percentage >= 70
        
        conn = get_db()
        cursor = conn.cursor()
        cursor.execute('''
            INSERT INTO quiz_attempts (quiz_id, student_name, score, total_questions, percentage, passed, answers, completed_at)
            VALUES (?, ?, ?, ?, ?, ?, ?, ?)
        ''', (
            quiz_id,
            student_name,
            score,
            total,
            percentage,
            passed,
            json.dumps(user_answers),
            datetime.now(timezone.utc).isoformat()
        ))
        conn.commit()
        conn.close()
        
        return jsonify({
            'success': True,
            'score': score,
            'total': total,
            'percentage': percentage,
            'passed': passed,
            'results': results
        }), 200
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/api/quiz/upload', methods=['POST'])
def upload_quiz():
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file uploaded'}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400
        
        if not allowed_file(file.filename):
            return jsonify({'error': 'File type not allowed'}), 400
        
        title = request.form.get('title', file.filename.rsplit('.', 1)[0])
        course_code = request.form.get('course_code', '').upper()
        description = request.form.get('description', '')
        difficulty = request.form.get('difficulty', 'medium')
        tags_input = request.form.get('tags', '')
        tags = [t.strip() for t in tags_input.split(',') if t.strip()]
        collection = request.form.get('collection', '').strip()
        
        try:
            num_questions = int(request.form.get('num_questions', 10))
        except:
            num_questions = 10
        
        filename = secure_filename(file.filename)
        os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(file_path)
        
        content = ""
        if filename.endswith('.pdf'):
            content = extract_text_from_pdf(file_path)
        elif filename.endswith(('.pptx', '.ppt')):
            content = extract_text_from_pptx(file_path)
        elif filename.endswith('.txt'):
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
        else:
            return jsonify({'error': 'Could not extract text from this file type'}), 400
        
        print(f"📄 Extracted {len(content)} characters for quiz")
        
        questions = ai_processor.generate_quiz_questions(content, num_questions, difficulty)
        
        quiz_id = str(uuid.uuid4())[:8]
        
        conn = get_db()
        cursor = conn.cursor()
        
        cursor.execute('''
            INSERT INTO quizzes (
                quiz_id, title, course_code, description, difficulty, questions,
                question_count, file_name, file_url, tags, upload_date, 
                views, downloads, rating, rating_count, collection
            ) VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
        ''', (
            quiz_id,
            title,
            course_code,
            description,
            difficulty,
            json.dumps(questions),
            len(questions),
            filename,
            f"/uploads/{filename}",
            json.dumps(tags),
            datetime.now(timezone.utc).isoformat(),
            0, 0, 0.0, 0,
            collection
        ))
        conn.commit()
        conn.close()
        
        return jsonify({
            'success': True,
            'message': f'Quiz created with {len(questions)} questions!',
            'quiz': {
                'quiz_id': quiz_id,
                'title': title,
                'course_code': course_code,
                'description': description,
                'difficulty': difficulty,
                'question_count': len(questions),
                'tags': tags,
                'collection': collection
            }
        }), 201
        
    except Exception as e:
        print(f"❌ Quiz upload error: {e}")
        return jsonify({'error': str(e)}), 500

@app.route('/api/quizzes/<quiz_id>/rate', methods=['POST'])
def rate_quiz(quiz_id):
    try:
        data = request.get_json()
        rating = data.get('rating', 0)
        
        if not 1 <= rating <= 5:
            return jsonify({'error': 'Rating must be between 1 and 5'}), 400
        
        conn = get_db()
        cursor = conn.cursor()
        cursor.execute('SELECT rating, rating_count FROM quizzes WHERE quiz_id = ?', (quiz_id,))
        row = cursor.fetchone()
        
        if not row:
            conn.close()
            return jsonify({'error': 'Quiz not found'}), 404
        
        current_rating = row['rating'] or 0
        current_count = row['rating_count'] or 0
        
        new_count = current_count + 1
        new_rating = round(((current_rating * current_count) + rating) / new_count, 1)
        
        cursor.execute('''
            UPDATE quizzes SET rating = ?, rating_count = ? WHERE quiz_id = ?
        ''', (new_rating, new_count, quiz_id))
        
        conn.commit()
        conn.close()
        
        return jsonify({
            'success': True,
            'rating': new_rating,
            'rating_count': new_count
        }), 200
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/uploads/<filename>')
def uploaded_file(filename):
    return send_from_directory(app.config['UPLOAD_FOLDER'], filename)

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in app.config['ALLOWED_EXTENSIONS']

if __name__ == '__main__':
    print("\n🚀 Starting USP Notes Marketplace API...")
    print("📡 Server running on: http://localhost:5000")
    print("💾 Database: SQLite")
    
    if ai_processor.provider:
        print(f"🤖 AI Model: {ai_processor.provider.upper()} - {ai_processor.model}")
    else:
        print("⚠️ No AI Provider - Using fallback processing")
    
    app.run(debug=True, port=5000)